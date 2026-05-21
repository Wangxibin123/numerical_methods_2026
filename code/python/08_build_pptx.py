r"""
08_build_pptx.py — Beamer → PNG → PPTX 自动化。

Beamer 输出 PDF，pdftoppm 转 PNG，本脚本把这些 PNG 全屏铺到一个新的
.pptx 文件里，并且把对应的中文讲稿（slides/narration.md 的每 § 段）
写到讲者备注 (notes pane)。

这样团队成员每次只需：
    1. 改 slides/main.tex     （Beamer 内容）
    2. 改 slides/narration.md （讲稿）
    3. make pptx              （自动重生成 ppt/lecture.pptx）
    4. 在 PowerPoint 中打开 ppt/lecture.pptx → Slide Show → Record Slide Show

录音和摄像头叠加由 PowerPoint 原生支持；我们这里只负责
自动同步 Beamer 内容 + 中文讲稿到 PPTX。

依赖：pip install python-pptx
"""

from __future__ import annotations

import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu, Inches

sys.path.insert(0, str(Path(__file__).resolve().parent))
import _config as cfg  # noqa: E402

SLIDES_DIR = cfg.REPO_ROOT / "slides"
RENDERED_DIR = SLIDES_DIR / "rendered"
NARRATION_MD = SLIDES_DIR / "narration.md"
PPTX_OUT = cfg.REPO_ROOT / "ppt" / "lecture.pptx"

# 16:9 widescreen, EMU = English Metric Units (914400 per inch)
SLIDE_WIDTH_EMU = Inches(13.333)
SLIDE_HEIGHT_EMU = Inches(7.5)


def parse_narration(md_path: Path) -> dict[int, dict[str, str]]:
    """Return {slide_index: {'title': ..., 'speaker': ..., 'script': ...}}.

    Each section header in narration.md looks like:
        ## §<n> <title> — <speaker>（约 N 秒）
    The body until the next §-header becomes the speaker note for that slide.
    """
    if not md_path.exists():
        return {}
    text = md_path.read_text(encoding="utf-8")
    header_re = re.compile(
        r"^##\s+§(?P<n>\d+)\s+(?P<title>[^—\-—]+?)\s*[—\-—]\s*"
        r"(?P<speaker>[^（(]+?)(?:[（(]\s*约\s*\S+\s*秒\s*[）)])?\s*$",
        re.MULTILINE,
    )

    out: dict[int, dict[str, str]] = {}
    matches = list(header_re.finditer(text))
    for k, m in enumerate(matches):
        start = m.end()
        end = matches[k + 1].start() if k + 1 < len(matches) else len(text)
        body = text[start:end].strip()
        # strip leading "> " quoting if author used quote-style narration
        lines = []
        for ln in body.splitlines():
            lines.append(re.sub(r"^>\s?", "", ln))
        body_clean = "\n".join(lines).strip()
        out[int(m.group("n"))] = {
            "title": m.group("title").strip(),
            "speaker": m.group("speaker").strip(),
            "script": body_clean,
        }
    return out


def collect_pngs() -> list[Path]:
    if not RENDERED_DIR.exists():
        return []
    pngs = sorted(RENDERED_DIR.glob("slide-*.png"))
    if not pngs:
        pngs = sorted(RENDERED_DIR.glob("slide*.png"))
    return pngs


def build_pptx(pngs: list[Path], narrations: dict[int, dict[str, str]]) -> Path:
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH_EMU
    prs.slide_height = SLIDE_HEIGHT_EMU

    # blank layout (index 6 in default template)
    blank_layout = prs.slide_layouts[6]

    for i, png in enumerate(pngs, start=1):
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(png), Emu(0), Emu(0),
            width=SLIDE_WIDTH_EMU, height=SLIDE_HEIGHT_EMU,
        )
        # speaker notes
        narration = narrations.get(i)
        notes_tf = slide.notes_slide.notes_text_frame
        if narration:
            notes_tf.text = (
                f"【{narration['title']}】录制：{narration['speaker']}\n\n"
                f"{narration['script']}"
            )
        else:
            notes_tf.text = f"（第 {i} 页未在 narration.md 找到对应讲稿）"

    PPTX_OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(PPTX_OUT))
    return PPTX_OUT


def main() -> int:
    pngs = collect_pngs()
    if not pngs:
        print(f"[err] no PNGs in {RENDERED_DIR}")
        print("  先跑：")
        print("    make slides     # 生成 slides/main.pdf")
        print("    make video-prep # 把 PDF 转成 PNG")
        return 1

    narrations = parse_narration(NARRATION_MD)
    if not narrations:
        print(f"[warn] no narration parsed from {NARRATION_MD}")
    else:
        print(f"[ok] parsed {len(narrations)} narration sections")
        missing = [i for i in range(1, len(pngs) + 1) if i not in narrations]
        if missing:
            print(f"[warn] slides without narration: {missing}")
        extra = [i for i in narrations if i > len(pngs)]
        if extra:
            print(f"[warn] narrations beyond available slides: {extra}")

    out = build_pptx(pngs, narrations)
    print(f"\nwrote {out}")
    print(f"  {len(pngs)} slides, {sum(1 for n in narrations.values() if n['script'])} with notes")
    print()
    print("下一步：")
    print(f"  1. 用 PowerPoint 打开 {out}")
    print("  2. View → Notes 检查每页讲稿")
    print("  3. Slide Show → Record Slide Show → 按讲稿录音 + 摄像头")
    print("  4. File → Export → Create Video → 1080p → final_video.mp4")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
